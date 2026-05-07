#!/usr/bin/env python3
"""Render each slide of a pptx into a single composited PNG.

Why: the dashboard is fully local; we don't want to depend on LibreOffice (700MB)
just to display review sheets. The review-sheet pptx files are mostly composed of
images positioned on a slide. We use python-pptx to walk shapes, get their EMU
positioning, and composite them onto a white canvas with PIL.

Caveats: text frames with custom fonts/effects don't render through this path —
they are written as faint placeholder text. For these particular review sheets
the content is overwhelmingly raster images so it works well.

Usage:
    python scripts/render_pptx.py <input.pptx> <output_dir> [--width 1600]
"""
from __future__ import annotations

import argparse
import io
import json
import pathlib
import sys
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes_recursive(shapes) -> Iterable:
    """Walk into group shapes recursively so we don't miss children."""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_recursive(sh.shapes)
        else:
            yield sh


def render_pptx(src: pathlib.Path, out_dir: pathlib.Path, target_w: int = 1600) -> dict:
    pres = Presentation(str(src))
    slide_w_emu = pres.slide_width
    slide_h_emu = pres.slide_height
    scale = target_w / slide_w_emu
    target_h = int(slide_h_emu * scale)

    out_dir.mkdir(parents=True, exist_ok=True)
    manifest = {
        "source": str(src),
        "slide_width": target_w,
        "slide_height": target_h,
        "slide_count": len(pres.slides),
        "slides": [],
    }

    # Try to use a basic font for text fallbacks
    font = None
    for candidate in [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]:
        if pathlib.Path(candidate).exists():
            try:
                font = ImageFont.truetype(candidate, size=22)
                break
            except Exception:
                pass
    if font is None:
        font = ImageFont.load_default()

    for idx, slide in enumerate(pres.slides):
        canvas = Image.new("RGB", (target_w, target_h), "white")
        first_text = None
        pic_count = 0
        for sh in iter_shapes_recursive(slide.shapes):
            try:
                left, top, width, height = sh.left, sh.top, sh.width, sh.height
            except Exception:
                continue
            if None in (left, top, width, height):
                continue
            x = max(0, int(left * scale))
            y = max(0, int(top * scale))
            w = max(1, int(width * scale))
            h = max(1, int(height * scale))

            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = Image.open(io.BytesIO(sh.image.blob)).convert("RGBA")
                    img = img.resize((w, h), Image.LANCZOS)
                    canvas.paste(img, (x, y), img if img.mode == "RGBA" else None)
                    pic_count += 1
                except Exception as e:
                    print(f"  slide {idx+1}: image error: {e}", file=sys.stderr)
            elif sh.has_text_frame:
                txt = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text.strip())
                if txt.strip():
                    if not first_text:
                        first_text = txt.strip().splitlines()[0][:80]
                    draw = ImageDraw.Draw(canvas)
                    # crude wrap
                    max_chars = max(20, w // 12)
                    lines = []
                    for raw_line in txt.splitlines():
                        while len(raw_line) > max_chars:
                            lines.append(raw_line[:max_chars])
                            raw_line = raw_line[max_chars:]
                        lines.append(raw_line)
                    line_h = 26
                    for i, line in enumerate(lines):
                        ly = y + i * line_h
                        if ly + line_h > y + h: break
                        draw.text((x + 6, ly + 4), line, fill=(40, 40, 40), font=font)

        out_file = out_dir / f"slide-{idx+1:02d}.png"
        canvas.save(out_file, "PNG", optimize=True)
        manifest["slides"].append({
            "index": idx + 1,
            "image": out_file.name,
            "title": first_text,
            "pictures": pic_count,
        })
        print(f"  slide {idx+1}/{len(pres.slides)}: {pic_count} pictures, title={first_text!r}")

    (out_dir / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return manifest


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("input", help="path to .pptx")
    ap.add_argument("out_dir", help="output directory for slide PNGs + manifest.json")
    ap.add_argument("--width", type=int, default=1600, help="target image width in pixels")
    args = ap.parse_args()
    src = pathlib.Path(args.input)
    if not src.exists():
        sys.exit(f"input not found: {src}")
    print(f"rendering {src.name} → {args.out_dir} @ {args.width}px wide")
    m = render_pptx(src, pathlib.Path(args.out_dir), target_w=args.width)
    print(f"done: {m['slide_count']} slides, manifest at {args.out_dir}/manifest.json")


if __name__ == "__main__":
    main()
